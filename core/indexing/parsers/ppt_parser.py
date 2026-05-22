"""ppt_parser.py — parse .ppt / .pptx files via python-pptx."""

import hashlib
from datetime import datetime, timezone
from pathlib import Path

from .base import BaseParser


def _md5(text: str) -> str:
    return hashlib.md5(text.encode("utf-8")).hexdigest()


class PptParser(BaseParser):
    """Parse PowerPoint files (both .ppt and .pptx) extracting text from shapes.

    Notes
    -----
    * Legacy .ppt files require LibreOffice to be installed on the system.
      The parser will attempt conversion using `libreoffice --headless`.
    * .pptx files are handled natively by python-pptx.
    """

    def parse(self, filepath: Path) -> dict:
        ext = filepath.suffix.lower()

        if ext == ".pptx":
            return self._parse_pptx(filepath)
        else:
            return self._parse_legacy(filepath)

    # ------------------------------------------------------------------
    def _parse_pptx(self, filepath: Path) -> dict:
        from pptx import Presentation

        prs = Presentation(str(filepath))
        slides_text: list[str] = []
        tables: list[list[str]] = []

        for slide in prs.slides:
            slide_lines: list[str] = []

            for shape in slide.shapes:
                # ---- text frames ----
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            slide_lines.append(line)

                # ---- tables ----
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            tables.append(cells)

            if slide_lines:
                slides_text.append("\n".join(slide_lines))

        content = "\n\n".join(slides_text)

        return {
            "source_file": filepath.name,
            "source_path": str(filepath),
            "content": content,
            "source_type": "pptx",
            "char_count": len(content),
            "line_count": content.count("\n") + 1,
            "md5": _md5(content),
            "parsed_at": datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"),
            "tables": tables,
        }

    # ------------------------------------------------------------------
    def _parse_legacy(self, filepath: Path) -> dict:
        """Convert .ppt to .pptx via LibreOffice headless, then parse."""
        import subprocess
        import tempfile
        import os

        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = os.path.join(tmpdir, filepath.name + ".pptx")
            try:
                subprocess.run(
                    [
                        "libreoffice",
                        "--headless",
                        "--convert-to", "pptx",
                        "--outdir", tmpdir,
                        str(filepath),
                    ],
                    check=True,
                    capture_output=True,
                    timeout=120,
                )
            except (subprocess.CalledProcessError, FileNotFoundError) as exc:
                raise RuntimeError(
                    "LibreOffice is required to parse .ppt files. "
                    "Install it from https://www.libreoffice.org/"
                ) from exc

            # LibreOffice may name the output file differently; find it
            converted = None
            for f in os.listdir(tmpdir):
                if f.lower().endswith(".pptx"):
                    converted = os.path.join(tmpdir, f)
                    break

            if not converted:
                raise RuntimeError(
                    f"LibreOffice conversion produced no .pptx for {filepath.name}"
                )

            result = self._parse_pptx(Path(converted))
            result["source_file"] = filepath.name
            result["source_path"] = str(filepath)
            result["source_type"] = "ppt"
            return result
